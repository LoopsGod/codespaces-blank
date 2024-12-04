import base64
import json
import logging
import os
import random
import re
import shutil
from io import BytesIO
from typing import List

from dotenv import load_dotenv
import fitz
import openai
import requests
from PIL import Image
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Inches, Pt
from pydantic import BaseModel

from spire.presentation import FileFormat, Presentation as Spire_Presentation
from spire.presentation.common import *


def get_slide_data(slide_id: str):
    """
    Function that given a slide id, returns the slide data.
    """

    SUPABASE_URL = "https://eckipwtivbjkzlwhtido.supabase.co"
    SUPABASE_SECRET_KEY = os.getenv('SUPABASE_SECRET_KEY')

    headers = {
        "prefer": "return=representation",
        "authorization": f"Bearer {SUPABASE_SECRET_KEY}",
        "content-type": "application/json"
    }

    url = f"{SUPABASE_URL}/rest/v1/slides?id=eq.{slide_id}&apikey={SUPABASE_SECRET_KEY}"

    response = requests.get(url, headers=headers)

    if response.ok:
        print("Data fetched successfully")
    else:
        print("Error fetching data:", response.text)
        return 400 
    return response.json()

def slides_count_to_int(slides_count: str):
    """
    function that given the slides_count variable, returns an integer
    if it is not a number such as "Automatic based on content"
    it returns a random number between 6 and 12.
    """
    if re.search(r'\D', str(slides_count)):
        return random.randint(6, 10)  
    return int(slides_count)

def download_file(slide_data_response):
    """ 
    function that gets slide_data_response and downloads the file (currently only support one file)
    saves the file in the tmp directory, returns None if no file and filepath if successful.
    """
    try:
        remote_url = slide_data_response[0].get("files")[0].get("remote")
        response = requests.get(remote_url)
        if response.status_code == 200:
            file_extension = slide_data_response[0].get("files")[0].get("name").split('.')[-1]
            tmp_dir = "/tmp"
            if not os.path.exists(tmp_dir):
                os.makedirs(tmp_dir)
            file_path = os.path.join(tmp_dir, f"uploaded_file.{file_extension}")
            with open(file_path, "wb") as file:
                file.write(response.content)
            print(f"File downloaded successfully as {file_path}")
            return file_path
        else:
            print("Failed to download file.")
            return None
    except Exception as e:
        print(f"No remote url found in the response: {e}")
        return None
    
def extract_text(input_file, max_words=4000):
    """
    Extracts text from a PDF or DOCX file.

    Parameters:
        input_file (str): Path to the PDF or DOCX file.

    Returns:
        str: Extracted text from the input file.
    """

    text = ""
    
    if input_file is None:
        return ""
    
    if input_file.endswith('.pdf'):
        try:
            with fitz.open(input_file) as doc:
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    text += page.get_text()
                    if len(re.findall(r'\w+', text)) >= max_words:
                        break
        except Exception as e:
            print(f"Error extracting text from PDF: {e}")
    elif input_file.endswith('.docx'):
        try:
            doc = Document(input_file)
            for para in doc.paragraphs:
                text += para.text + '\n'
                if len(re.findall(r'\w+', text)) >= max_words:
                    break
        except Exception as e:
            print(f"Error extracting text from DOCX: {e}")
    elif isinstance(input_file, str):
        text = input_file[:max_words]
        return text

    if not isinstance(text, str):
        raise ValueError("Unsupported input type.")
    

    return text 

def extract_images(input_data, output_dir, max_images=10):
    """
    Extracts images from a PDF file, DOCX file, or a base64-encoded string.

    Parameters:
        input_data (str): Path to the PDF/DOCX file or the base64 string.
        output_dir (str): Directory where extracted images will be saved.
        max_images (int): Maximum number of images to extract (default is 10).

    Raises:
        ValueError: If the input type is unsupported or if extraction fails.
    """

    os.makedirs(output_dir, exist_ok=True)
    image_count = 0

    def save_and_verify(image_bytes, image_filename):
        nonlocal image_count
        if image_count >= max_images:
            return False
        image_path = os.path.join(output_dir, image_filename)
        with open(image_path, 'wb') as img_file:
            img_file.write(image_bytes)
        try:
            with Image.open(image_path) as img:
                img.verify()
            print(f"Saved image: {image_filename}")
            image_count += 1
            if image_count >= max_images:
                print(f"Image extraction limit of {max_images} reached.")
                return False
            return True
        except Exception as e:
            print(f"Invalid image {image_filename}: {e}")
            os.remove(image_path)
            return True

    if isinstance(input_data, str):
        if os.path.isfile(input_data):
            file_ext = os.path.splitext(input_data)[1].lower()
            if file_ext == '.pdf':
                input_type = 'pdf'
            elif file_ext == '.docx':
                input_type = 'docx'
            else:
                raise ValueError("Unsupported file extension. Supported extensions are '.pdf' and '.docx'.")
        else:
            input_type = 'string'
    else:
        raise ValueError("input_data must be a file path or a base64-encoded string.")

    if input_type == 'pdf':
        try:
            doc = fitz.open(input_data)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list, start=1):
                    if image_count >= max_images:
                        break
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    image_filename = f'page_{page_num + 1}_img_{img_index}.{image_ext}'
                    should_continue = save_and_verify(image_bytes, image_filename)
                    if not should_continue:
                        break
                if image_count >= max_images:
                    break
            doc.close()
        except Exception as e:
            print(f"Error processing PDF: {e}")

    elif input_type == 'docx':
        try:
            doc = Document(input_data)
            rels = doc.part._rels
            img_count = 1
            for rel in rels:
                if image_count >= max_images:
                    break
                rel_obj = rels[rel]
                if "image" in rel_obj.target_ref:
                    image_bytes = rel_obj.target_part.blob
                    image_ext = os.path.splitext(rel_obj.target_ref)[1].replace('.', '').lower()
                    if image_ext not in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff']:
                        image_ext = 'png'
                    image_filename = f'image_{img_count}.{image_ext}'
                    should_continue = save_and_verify(image_bytes, image_filename)
                    if not should_continue:
                        break
                    img_count += 1
        except Exception as e:
            print(f"Error processing DOCX: {e}")

    elif input_type == 'string':
        try:
            if image_count < max_images:
                image_data = base64.b64decode(input_data)
                image = Image.open(BytesIO(image_data))
                image.verify()
                img_format = image.format.lower()
                image_filename = f'extracted_image.{img_format}'
                save_and_verify(image_data, image_filename)
        except Exception as e:
            print(f"Error processing string input: {e}")

    else:
        raise ValueError("Unsupported input type.")

    print(f"Total images extracted and saved: {image_count}")

def get_template_category_meta_data(template_category):
    """
    Function that given a template category, returns the meta data for that category.
    """
    
    url = f"https://cdn.qlina.ai/templates/{template_category}/_meta.json"
    print(url)
    response = requests.get(url)
    if response.ok:
        return response.json()
    else:
        return None

def presentation_picker_from_meta_data(text_content, prompt, tone, meta_data, template_type):
    """
    Function that, given the text content, prompt, tone, and meta data, returns a template - response in json format.
    It saves the chosen template in tmp as template.pptx.
    """

    client = openai.OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
    
    class TemplateSelection(BaseModel):
        selected_template: str

    messages = [
        {
            "role": "system",
            "content": "You are an expert in selecting appropriate presentation templates based on content and tone.",
        },
        {
            "role": "user",
            "content": f"""
        Based on the following information, select the most appropriate .pptx template for a presentation.

        ### Available Templates:
        {meta_data}

        ### Presentation Content:
        {text_content}

        ### Prompt:
        {prompt}

        ### Tone:
        {tone}

        ### Instructions:
        From the available templates, select the one that best matches the content and tone of the presentation.
        Only select one template from the list provided in the meta data.
        Return your answer in the following JSON format:

        {{
            "selected_template": "<template_name>"
        }}

        Replace <template_name> with the name of the selected template.
        Ensure that the template name is exactly as given in the meta data.

        Do not include any additional information or explanation.
        """,
                },
            ]

    try:
        completion = client.beta.chat.completions.parse(
            model="gpt-4o-mini",
            messages=messages,
            response_format=TemplateSelection,
            max_tokens=100,
            temperature=0.3,
        )

        message = completion.choices[0].message
        # print(completion)

        if message.refusal:
            print("The model refused to select a template.")
            return None
        elif message.parsed:
            template_selection = message.parsed
            print(f"Selected template: {template_selection.selected_template}")

             # download and store chosen template to tmp/template.pptx
            try:
                template_url = f"https://qlina.b-cdn.net/templates/{template_type}/{template_selection.selected_template}"
                response = requests.get(template_url)
                response.raise_for_status()
                with open("/tmp/template.pptx", "wb") as file:
                    file.write(response.content)
                return template_selection
            except Exception as e:
                print(f"An error occurred while downloading the template: {e}")
                return None
        else:
            print("The model did not return a valid template selection.")
            return None

    except Exception as e:
        print(f"An error occurred: {e}")
        return None

def parse_template_slides(pptx_file):
    """
    Function that parses a PowerPoint presentation and extracts information about each slide,
    it returns a list of dictionaries containing information about each slide.
    """
    prs = Presentation(pptx_file)
    slides_info = []

    def process_shape(shape, slide_info):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for shp in shape.shapes:
                process_shape(shp, slide_info)
        else:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_info['images'].append({
                    'name': shape.name
                })
            elif shape.has_text_frame and shape.text.strip():
                slide_info['text_boxes'].append({
                    'position': (shape.left, shape.top),
                    'size': (shape.width, shape.height),
                    'placeholder_text': shape.text.strip()
                })

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_info = {
            'template_slide_number': slide_number,
            'text_boxes': [],
            'images': [],
            'notes': ''
        }

        for shape in slide.shapes:
            process_shape(shape, slide_info)

        slides_info.append(slide_info)

    return slides_info

def generate_presentation_plan(essay_text, template_structure, total_slides):
    """ 
    Function that generates a presentation plan based on the essay text, template structure, and total number of slides.
    """

    client = openai.OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

    class SlideSelection(BaseModel):
        template_slide_number: int

    class PresentationPlan(BaseModel):
        opening_slide: SlideSelection
        content_slides: List[SlideSelection]
        conclusion_slide: SlideSelection

    messages = [
        {"role": "system", "content": "You are an expert in making and structuring presentations."},
        {"role": "user", "content": f"""
        You are an assistant that helps find the right order for presentation slides based on a template essay, where each slide can be used as a building block for our own presentation.
        Below is the template slide information extracted from a PowerPoint file, followed by the essay text on which the presentation will be later based, first however we have to make the structure.
        Ensure that the presentation always starts with an opening slide and ends with a conclusion slide. Organize the content slides in between based on the essay's structure and content.

        \"\"\"
        ### Template Slides:
        {template_structure}
        \"\"\"

        And the essay text:

        \"\"\"
        ### Essay Text:
        {essay_text}
        \"\"\"


         ### Desired Number of Slides:
        {total_slides}


        ### Instructions:
        Based on the essay text and the number of slides, which is {total_slides}, generate a JSON structure that orders the slides appropriately.
        The first slide must be the opening slide, the last slide must be the conclusion slide, and the content slides should be ordered logically based on the essay's content.
        Use the `template_slide_number` from the template slides to reference which slides to use.

        ### JSON Structure:
        {{
            "presentation_plan": {{
                "opening_slide": {{
                    "template_slide_number": <int>
                }},
                "content_slides": [
                    {{
                        "template_slide_number": <int>
                    }},
                    // ... more content slides
                ],
                "conclusion_slide": {{
                    "template_slide_number": <int>
                }}
            }}
        }}

        Ensure that the output adheres strictly to the provided schema.
        AND DO EXACTLY {total_slides}. YOU SOMETIMES DO TOO LITTLE
        """}
            ]

    try:
        completion = client.beta.chat.completions.parse(
            model="gpt-4o-mini",
            messages=messages,
            response_format=PresentationPlan,
            max_tokens=2000,
            temperature=0.1,
        )

        message = completion.choices[0].message
        # print(completion)

        if message.refusal:
            print("The model refused to generate the presentation plan.")
            return None
        elif message.parsed:
            presentation_plan = message.parsed
            return presentation_plan
        else:
            print("The model refused to generate the presentation plan.")
            return None

    except Exception as e:
        print(f"An error occurred: {e}")
        return None

def get_slide_placeholders(presentation_plan, template_slides):
    """
    Function that extracts placeholders from each slide in the presentation plan based on the template slides.
    Such that the prompt can fill in the output in a pre-determined structure.
    """
    slides_content = []

    def extract_placeholders(slide_number):
        slide_info = template_slides[slide_number - 1]
        placeholders = []
        for idx, tb in enumerate(slide_info['text_boxes']):
            placeholder_text = tb['placeholder_text']
            placeholders.append({
                'placeholder_id': f"{idx+1}",
                'position': {
                    'left': tb['position'][0],
                    'top': tb['position'][1]
                },
                # 'size': {
                #     'width': tb['size'][0],
                #     'height': tb['size'][1]
                # },
                'placeholder_text': placeholder_text,
                'max_char_length': len(placeholder_text)
            })
        return placeholders

    opening_slide_number = presentation_plan.opening_slide.template_slide_number
    opening_placeholders = extract_placeholders(opening_slide_number)
    slides_content.append({
        'template_slide_number': opening_slide_number,
        'placeholders': opening_placeholders
    })

    for content_slide in presentation_plan.content_slides:
        slide_number = content_slide.template_slide_number
        placeholders = extract_placeholders(slide_number)
        slides_content.append({
            'template_slide_number': slide_number,
            'placeholders': placeholders
        })

    conclusion_slide_number = presentation_plan.conclusion_slide.template_slide_number
    conclusion_placeholders = extract_placeholders(conclusion_slide_number)
    slides_content.append({
        'template_slide_number': conclusion_slide_number,
        'placeholders': conclusion_placeholders
    })

    return slides_content

def generate_presentation_content(essay_text, prompt, tone, slides_content):
    """
    Function that generates presentation content based on the total text, prompt, tone and slide placeholders.
    """

    client = openai.OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

    class SlideContent(BaseModel):
        template_slide_number: int
        contents: List[str]

    class PresentationContentResponse(BaseModel):
        slides: List[SlideContent]

    messages = [
            {
                "role": "system",
                "content": (
                    "You are an expert at creating presentation content based on provided placeholders. "
                    "Each placeholder includes initial text, position, size, a unique identifier, and a maximum character length. "
                    "Your task is to generate replacement content that fits exactly within these placeholders without exceeding the maximum character length."
                )
            },
            {
                "role": "user",
                "content": f"""
        Based on the essay text and the slide placeholders provided, generate the content for each slide. For each slide, replace the placeholder text with appropriate content derived from the essay.

        \"\"\"
        ### Essay Text:
        {essay_text} + with additional request: {prompt} + in a {tone} tone
        \"\"\"

        \"\"\"
        ### Slides and Placeholders:
        {json.dumps(slides_content, indent=4)}
        \"\"\"

        ### Instructions:
        - **Replacement Constraints**: For each placeholder, generate concise and relevant content that **does not exceed** the `max_char_length`. The replacement should fit within the space provided.
        - **Placeholder Identification**: Use the `placeholder_id` to ensure that content is placed in the correct placeholder. **Do not swap** content between placeholders.
        - **Size and Position**: Consider the position and size of each placeholder to ensure the content fits appropriately.
        - **Language Consistency**: If the essay text is in a language other than English, use that same language in your content.
        - **Relevance**: Ensure the content aligns with the overall theme and key points of the essay.
        - **Formatting**: Use bullet points if appropriate.
        - **Output Format**: Return the content in the following JSON format:

        {{
            "slides": [
                {{
                    "template_slide_number": <int>,
                    "contents": [
                        {{
                            "placeholder_id": "<unique_placeholder_id>",
                            "content": "<replacement content>"
                        }},
                        // ... more placeholders
                    ]
                }},
                // ... more slides
            ]
        }}

        ### Additional Guidelines:
        - **Text Length Constraint**: The length of the replacement content must **not exceed** `max_char_length`. If necessary, rephrase to maintain brevity.
        - **Avoid Placeholder Swapping**: Ensure that each piece of content replaces the correct placeholder based on its `placeholder_id`. Do not swap content between placeholders.
        - **Consistency and Accuracy**: Double-check that each replacement accurately reflects the intended content for its specific placeholder without introducing errors or irrelevant information.

        Ensure that the output adheres **strictly** to the provided schema and instructions.
        """
            }
        ]


    try:
        completion = client.beta.chat.completions.parse(
            model="gpt-4o-2024-08-06",
            messages=messages,
            response_format=PresentationContentResponse,
            max_tokens=4000,
            temperature=0.1,
        )

        if completion.choices[0].message.refusal:
            print("The model refused to generate the presentation content.")
            return None

        response_text = completion.choices[0].message.content.strip()
        presentation_content = json.loads(response_text)
        return presentation_content

    except Exception as e:
        print(f"An error occurred: {e}")
        return None
    
def add_image_info_to_presentation_content(presentation_content, template_slides):
    """
    Function that adds image information to the presentation content based on the original template presentation.
    """

    for slide_data in presentation_content['slides']:
        slide_number = slide_data['template_slide_number']
        slide_info = template_slides[slide_number - 1]

        images = []
        for shape in slide_info['images']:
            images.append({
                'name': shape['name']
            })

        slide_data['images'] = images

    return presentation_content

def encode_image(image_path):
    """
    Function that encodes an image file to base64.
    """ 
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')
    
def assign_images(image_folder, presentation_content):
    """
    Function that assigns images to presentation slides based on their content.
    Also adds a brief description for image retrieval is necessary
    """
    client = openai.OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

    class ImageAssignment(BaseModel):
        slide_number: int
        image_filename: str  # Use the image filename ONLY if you have assigned an image
        brief_description: str  # provide a brief description ONLY IF NO image has be assigned

    class ImageAssignmentResponse(BaseModel):
        assignments: List[ImageAssignment]

    images_needed = 0
    for slide in presentation_content['slides']:
        if 'images' in slide and len(slide['images']) > 0:
            images_needed += len(slide['images'])

    print(f"Number of images needed: {images_needed}")

    if images_needed == 0:
        print("No images needed for the presentation.")
        return

    # Check if images are available in 'file_images' folder
    image_files = [f for f in os.listdir(image_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'))]

    if not image_files:
        print("No images found in the folder.")
        return

    system_prompt = '''
        You are an AI assistant that helps assign images to presentation slides based on their content.

        **Instructions:**

        - You will be provided with a set of images along with their filenames.
        - You will also be provided with the content of each slide that requires an image.
        - Your task is to match each image to the most appropriate slide based on the image content and the slide content. Don't match if completely irrelevant
        - If the number of images is less than the number of slides requiring images, assign images to the most relevant slides.
        - For slides without an assigned image, write a brief description (4-5 words) of what the image should depict.

        **Output Format:**

        Return the result in **JSON format**:

        {
            "assignments": [
                {
                    "slide_number": <int>,
                    "image_filename": "<string>"  // Use the image filename ONLY if you have assigned an image
                    "brief_description": "<string>" // provide a brief description ONLY IF NO image has be assigned
                },
                // ...
            ]
        }

        Ensure that your response is valid JSON and matches the exact structure shown above.
    '''

    messages = [{"role": "system", "content": system_prompt}]

    user_content = []

    encoded_images = []
    for image_file in image_files:
        image_path = os.path.join(image_folder, image_file)
        base64_image = encode_image(image_path)
        encoded_images.append((image_file, base64_image))

    # Add images to the user content
    for image_filename, base64_image in encoded_images:
        user_content.extend([
            {
                "type": "text",
                "text": f"Image: {image_filename}"
            },
            {
                "type": "image_url",
                "image_url": {
                    "url": f"data:image/jpeg;base64,{base64_image}",
                    "detail": "low"
                }
            }
        ])

    # Add presentation_content slides information
    slides_info = []
    for idx, slide in enumerate(presentation_content['slides']):
        if 'images' in slide and len(slide['images']) > 0:
            slide_number = idx + 1
            slide_content_text = ' '.join(slide['contents'])
            slides_info.append(f"Slide {slide_number}: {slide_content_text}")


    slides_text = "\n".join(slides_info)
    user_content.append({
        "type": "text",
        "text": f"Slides Content:\n{slides_text}"
    })

    messages.append({
        "role": "user",
        "content": user_content
    })

    # print(messages)

    try:
        response = client.beta.chat.completions.parse(
            model="gpt-4o-mini",
            messages=messages,
            response_format=ImageAssignmentResponse,
            max_tokens=1000,
            temperature=0.3,
        )
    except Exception as e:
        print(f"An error occurred: {e}")
        return None

    if response.choices[0].message.refusal:
        print("The model refused to generate the image assignments.")
        return None
    elif response.choices[0].message.parsed:
        image_assignments = response.choices[0].message.parsed.assignments
        return image_assignments
    else:
        print("The model did not return valid assignments.")
        return None
    
def fetch_and_download_images(assignments, image_folder):
    """
    Function that fetches images from PIXABAY based on the brief description and 
    downloads them to the folder /tmp/images. It also changes the image_assigments variable to include the path.
    """

    PIXABAY_API_KEY = os.getenv('PIXABAY_API_KEY')

    if assignments is None:
        return None

    for assignment in assignments:
        if not assignment.image_filename:
            query = assignment.brief_description
            url = f"https://pixabay.com/api/?key={PIXABAY_API_KEY}&q={query}&image_type=photo&per_page=3"

            #  make the API request
            response = requests.get(url)
            if response.status_code == 200:
                data = response.json()
                if data['hits']:
                    image_url = data['hits'][0]['largeImageURL']
                    filename = f"slide_{assignment.slide_number}.jpg"
                    image_path = os.path.join(image_folder, filename)
                    print(image_path)

                    img_data = requests.get(image_url).content
                    with open(image_path, 'wb') as img_file:
                        img_file.write(img_data)

                    assignment.image_filename = filename
                    print(f"Downloaded and assigned image for slide {assignment.slide_number}: {filename}")
                else:
                    print(f"No images found for: {query}")
            else:
                print(f"Failed to fetch images for: {query} (Status code: {response.status_code})")

    return assignments

def to_list_structure_template(output_structure):
    """
    Function that converts the output_structure to a list structure. 
    For looping purposes. Used in the generate_presentation_from_plan function.
    """

    list_structure = []
    list_structure.append(output_structure.opening_slide.template_slide_number)
    for content_slide in output_structure.content_slides:
        list_structure.append(content_slide.template_slide_number)
    list_structure.append(output_structure.conclusion_slide.template_slide_number)

    return list_structure

def generate_presentation_from_plan(output_structure, template_file, empty_qlina):
    """
    Function that generates the presentation from the plan using Spire's copy functionality.
    """
    outputFile = '/tmp/ordered_template.pptx'

    sourcePPT = Spire_Presentation()
    sourcePPT.LoadFromFile(template_file)

    destPPT = Spire_Presentation()
    destPPT.LoadFromFile(empty_qlina)

    list_structure = to_list_structure_template(output_structure)

    for slide_number in list_structure:
        slide1 = sourcePPT.Slides[slide_number - 1]
        destPPT.Slides.AppendBySlide(slide1)


    destPPT.SaveToFile(outputFile, FileFormat.Pptx2016)
    destPPT.Dispose()
    return outputFile

def remove_evaluation_warnings(input_pptx, output_pptx):
    """
    Function that removes evaluation warnings from Spire PowerPoint presentation.
    """

    prs = Presentation(input_pptx)

    target_text = 'Evaluation Warning : The document was created with Spire.Presentation for Python'

    for slide_number, slide in enumerate(prs.slides, start=1):
        shapes_to_remove = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()

                if text == target_text:
                    shapes_to_remove.append(shape)

        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)

    prs.save(output_pptx)
    return output_pptx

def insert_content_into_presentation(presentation_file, presentation_content, save_as):
    """
    Function that inserts content into a pptx presentation based on the provided content.
    """
    prs = Presentation(presentation_file)
    total_slides = len(prs.slides)

    counter = 0
    for slide_data in presentation_content['slides']:
        slide_number = counter
        counter += 1
        contents = slide_data['contents']

        if 0 <= slide_number < total_slides:
            slide = prs.slides[slide_number]
        else:
            continue

        content_idx = 0

        def process_shape(shape):
            nonlocal content_idx
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for shp in shape.shapes:
                    process_shape(shp)
            elif shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if content_idx < len(contents):
                            old_text = run.text
                            new_text = contents[content_idx]
                            run.text = new_text
                            content_idx += 1
                        else:
                            break
                    if content_idx >= len(contents):
                        break
            if content_idx >= len(contents):
                return

        for shape in slide.shapes:
            process_shape(shape)
            if content_idx >= len(contents):
                break

    prs.save(save_as)

def replace_images_in_presentation(pptx_file, image_assignments, output_file, image_folder):
    """
    Replaces images on specific slides in a PowerPoint presentation, including images inside group shapes.

    Parameters:
        pptx_file (str): Path to the input PowerPoint file.
        image_assignments (List[ImageAssignment]): List of image assignments with slide numbers and image filenames.
        output_file (str): Path to save the modified PowerPoint file.
        image_folder (str): Path to the folder containing images.

    Returns:
        Nothing
    """
    prs = Presentation(pptx_file)

    if image_assignments is None:
        prs.save(output_file)
        return None

    # Build a mapping from slide numbers to image filenames
    image_replacements = {}
    for assignment in image_assignments:
        slide_number = assignment.slide_number
        image_filename = assignment.image_filename
        if image_filename:
            image_replacements.setdefault(slide_number, []).append(image_filename)

    for slide_number, image_filenames in image_replacements.items():
        if 1 <= slide_number <= len(prs.slides):
            slide = prs.slides[slide_number - 1]

            # Create an iterator over the image filenames for the current slide
            image_iterator = iter(image_filenames)

            # Function to process shapes recursively
            def process_shape(shape):
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for shp in shape.shapes:
                        process_shape(shp)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_filename = next(image_iterator)
                        image_path = os.path.join(image_folder, image_filename)
                        if os.path.exists(image_path):
                            try:
                                with Image.open(image_path) as img:
                                    img.verify()
                                # Remove the old image
                                sp = shape.element
                                sp.getparent().remove(sp)
                                # Add the new image at the same position and size
                                slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)
                            except Exception as e:
                                print(f"File is not a valid image: {image_path}, error: {e}")
                        else:
                            print(f"Image file not found: {image_path}")
                    except StopIteration:
                        pass
            for shape in slide.shapes:
                process_shape(shape)

    prs.save(output_file)
    print(f"Images replaced and saved to {output_file}")

def upload_final_to_bunny(file_path, slides_id):
    """
    Function that uploads the final presentation to BunnyCDN.
    """
    BUNNY_API_KEY = os.getenv("BUNNY_API_KEY")

    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"The file '{file_path}' does not exist.")

    upload_url = f"https://storage.bunnycdn.com/qlina/slides/{slides_id}.pptx"
    headers = {
        "AccessKey": BUNNY_API_KEY,
        "Content-Type": "application/octet-stream",
    }

    with open(file_path, 'rb') as file:
        binary_data = file.read()

    response = requests.put(upload_url, data=binary_data, headers=headers)

    if response.status_code in (200, 201):
        print(f"File uploaded successfully to bunny.")
    else:
        print(f"Failed to upload file. Status Code: {response.status_code}")
        print(f"Response: {response.text}")

    return response

def initaite_pptx_2_pdf(slides_id):
    url = "https://hkdk.events/hjz5e8xmigvami"
    input_url = f"https://qlina.b-cdn.net/slides/{slides_id}.pptx"
    payload = {
        "id": slides_id,
        "input_url": input_url,
        "input_file_name": f"{slides_id}.pptx",
        "export_file_name": f"{slides_id}.pdf",
        "event": "slides.done"
    }

    try:
        response = requests.post(url, json=payload)
        response.raise_for_status() 
        print("Response:", response.json())  
    except requests.exceptions.RequestException as e:
        print("An error occurred:", e)

def update_supabase(slides_id, title):
    """
    Function that given a slides_id and title, updates the supabase.
    """

    SUPABASE_URL = "https://eckipwtivbjkzlwhtido.supabase.co"
    SUPABASE_SECRET_KEY = os.getenv('SUPABASE_SECRET_KEY')

    headers = {
        "prefer": "return=representation",
        "authorization": f"Bearer {SUPABASE_SECRET_KEY}",
        "content-type": "application/json"
    }

    url = f"{SUPABASE_URL}/rest/v1/slides?id=eq.{slides_id}&apikey={SUPABASE_SECRET_KEY}"


    update_fields = {
        "status": "done",
        "title": title,
        "pptx_download_url": f"https://qlina.b-cdn.net/templates/slides/{slides_id}.pptx",
        "pdf_download_url": f"https://qlina.b-cdn.net/templates/slides/{slides_id}.pdf",
    }

    response = requests.patch(url, headers=headers, json=update_fields)

    if response.ok:
        print("Data updated successfully")
        return response.json()
    else:
        print("Error updating data:", response.text)
        return 400
